import os
import re
import uuid
import json
import socket
from io import BytesIO
from pathlib import Path
from datetime import datetime, timedelta, timezone
from zoneinfo import ZoneInfo
from typing import List, Literal

import httpx
from fastapi import FastAPI, Form, HTTPException, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field
from pypdf import PdfReader
from pptx import Presentation

OLLAMA_BASE_URL = os.getenv("OLLAMA_BASE_URL", "http://127.0.0.1:11434").rstrip("/")
DEFAULT_MODEL = "gemma3:latest"
VISION_MODEL_HINTS = (
    "gemma3",
    "llava",
    "bakllava",
    "llama3.2-vision",
    "moondream",
    "qwen2-vl",
    "minicpm-v",
)
MAX_UPLOAD_BYTES = 50 * 1024 * 1024
MAX_CHUNKS_PER_SESSION = 300
RAG_STORE_PATH = Path(os.getenv("RAG_STORE_PATH", "backend/rag_store.json"))

app = FastAPI(title="NOS LLM Chat API", version="2.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
)


class ChatMessage(BaseModel):
    role: Literal["system", "user", "assistant"]
    content: str = Field(min_length=1)


class ChatRequest(BaseModel):
    message: str = Field(min_length=1)
    model: str = DEFAULT_MODEL
    history: List[ChatMessage] = Field(default_factory=list)
    image_base64: str | None = None
    image_mime_type: str | None = None
    image_name: str | None = None
    rag_session_id: str | None = None
    use_rag: bool = True


class ChatResponse(BaseModel):
    response: str
    model: str


class RagUploadResponse(BaseModel):
    session_id: str
    doc_id: str
    file_name: str
    chunks: int


class RagSessionClearRequest(BaseModel):
    session_id: str


class RagRenameRequest(BaseModel):
    session_id: str
    doc_id: str | None = None
    doc_index: int | None = None
    old_file_name: str | None = None
    file_name: str = Field(min_length=1, max_length=200)


class RagDeleteRequest(BaseModel):
    session_id: str
    doc_id: str | None = None
    doc_index: int | None = None
    old_file_name: str | None = None


class RagDocumentUpdateRequest(BaseModel):
    session_id: str
    doc_id: str | None = None
    doc_index: int | None = None
    old_file_name: str | None = None
    content: str = Field(min_length=1)


RAG_STORE: dict[str, list[dict]] = {}


def _save_rag_store() -> None:
    try:
        payload: dict[str, list[dict]] = {}
        for session_id, docs in RAG_STORE.items():
            serialized_docs: list[dict] = []
            for doc in docs:
                serialized_docs.append(
                    {
                        "doc_id": doc.get("doc_id", str(uuid.uuid4())),
                        "file_name": doc.get("file_name", ""),
                        "chunks": doc.get("chunks", []),
                        # set is not JSON serializable -> persist as sorted list
                        "tokens": [sorted(list(t)) for t in doc.get("tokens", [])],
                    }
                )
            payload[session_id] = serialized_docs

        RAG_STORE_PATH.parent.mkdir(parents=True, exist_ok=True)
        RAG_STORE_PATH.write_text(json.dumps(payload, ensure_ascii=False), encoding="utf-8")
    except Exception:
        # Keep app available even if persistence fails.
        pass


def _load_rag_store() -> None:
    if not RAG_STORE_PATH.exists():
        return
    needs_persist = False
    try:
        raw = RAG_STORE_PATH.read_text(encoding="utf-8")
        data = json.loads(raw)
        if not isinstance(data, dict):
            return

        loaded: dict[str, list[dict]] = {}
        for session_id, docs in data.items():
            if not isinstance(session_id, str) or not isinstance(docs, list):
                continue
            parsed_docs: list[dict] = []
            for doc in docs:
                if not isinstance(doc, dict):
                    continue
                doc_id = doc.get("doc_id")
                file_name = doc.get("file_name")
                chunks = doc.get("chunks")
                tokens = doc.get("tokens")
                normalized_doc_id = str(doc_id).strip() if isinstance(doc_id, str) else ""
                if not normalized_doc_id:
                    normalized_doc_id = str(uuid.uuid4())
                    needs_persist = True
                if not isinstance(file_name, str) or not isinstance(chunks, list):
                    continue
                if not isinstance(tokens, list):
                    # Backward compatibility: rebuild tokens if missing
                    needs_persist = True
                    parsed_docs.append(
                        {
                            "doc_id": normalized_doc_id,
                            "file_name": file_name,
                            "chunks": chunks,
                            "tokens": [_tokenize(c) for c in chunks if isinstance(c, str)],
                        }
                    )
                    continue
                parsed_tokens = []
                for tok_list in tokens:
                    if isinstance(tok_list, list):
                        parsed_tokens.append(set(str(x) for x in tok_list))
                    else:
                        needs_persist = True
                        parsed_tokens.append(set())
                parsed_docs.append(
                    {
                        "doc_id": normalized_doc_id,
                        "file_name": file_name,
                        "chunks": [c for c in chunks if isinstance(c, str)],
                        "tokens": parsed_tokens,
                    }
                )
            if parsed_docs:
                loaded[session_id] = parsed_docs
        RAG_STORE.clear()
        RAG_STORE.update(loaded)
        if needs_persist:
            _save_rag_store()
    except Exception:
        # Ignore corrupted file and continue with empty in-memory store.
        pass


def _extract_ollama_error(exc: httpx.HTTPStatusError) -> str:
    try:
        payload = exc.response.json()
        if isinstance(payload, dict):
            if isinstance(payload.get("error"), str) and payload["error"].strip():
                return payload["error"]
            if isinstance(payload.get("detail"), str) and payload["detail"].strip():
                return payload["detail"]
    except ValueError:
        pass
    text = (exc.response.text or "").strip()
    return text if text else "Ollama API 오류가 발생했습니다."


async def ollama_request(method: str, path: str, payload: dict | None = None, timeout: float = 120.0):
    # Local LLM calls should ignore system proxy settings to avoid localhost routing issues.
    async with httpx.AsyncClient(timeout=timeout, trust_env=False) as client:
        try:
            response = await client.request(method, f"{OLLAMA_BASE_URL}{path}", json=payload)
            response.raise_for_status()
            return response.json()
        except httpx.ConnectError:
            raise HTTPException(
                status_code=503,
                detail="Ollama 서버에 연결할 수 없습니다. Ollama 실행 상태를 확인해주세요.",
            )
        except httpx.HTTPStatusError as exc:
            detail = _extract_ollama_error(exc)
            raise HTTPException(status_code=exc.response.status_code, detail=detail)
        except httpx.ReadTimeout:
            raise HTTPException(status_code=504, detail="Ollama 응답 시간이 초과되었습니다.")


def _history_to_prompt(history: List[ChatMessage], user_message: str) -> str:
    lines: List[str] = []
    for msg in history:
        role = "Assistant" if msg.role == "assistant" else ("System" if msg.role == "system" else "User")
        lines.append(f"{role}: {msg.content}")
    lines.append(f"User: {user_message}")
    lines.append("Assistant:")
    return "\n".join(lines)


def _current_date_system_message() -> str:
    try:
        now_kst = datetime.now(ZoneInfo("Asia/Seoul"))
    except Exception:
        # Fallback for Windows environments without tzdata.
        now_kst = datetime.now(timezone(timedelta(hours=9)))
    today_kr = f"{now_kst.year}년 {now_kst.month}월 {now_kst.day}일"
    # Keep this concise and deterministic so the model consistently answers date/time questions.
    return (
        "시스템 시간 기준 정보:\n"
        f"- 현재 한국 표준시(KST): {now_kst.strftime('%Y-%m-%d %H:%M:%S')}\n"
        f"- 오늘 날짜: {today_kr}\n"
        "사용자가 날짜/시간을 물으면 위 기준을 우선으로 답하세요."
    )


def _now_kst() -> datetime:
    try:
        return datetime.now(ZoneInfo("Asia/Seoul"))
    except Exception:
        return datetime.now(timezone(timedelta(hours=9)))


def _is_datetime_query(text: str) -> bool:
    t = (text or "").lower()
    keywords = [
        "오늘", "현재", "지금", "날짜", "시간", "몇시", "몇 시", "몇월", "몇 월", "며칠",
        "kst", "한국 시간", "한국시간",
        "today", "date", "time", "current time", "what time",
    ]
    return any(k in t for k in keywords)


def _format_datetime_answer() -> str:
    now = _now_kst()
    return (
        "현재 대한민국 표준시(KST) 기준:\n"
        f"- 날짜: {now.year}년 {now.month}월 {now.day}일\n"
        f"- 시간: {now.strftime('%H시 %M분 %S초')}\n"
        f"- ISO: {now.strftime('%Y-%m-%d %H:%M:%S')} (KST)"
    )


def _decode_text_file(raw: bytes) -> str:
    for encoding in ("utf-8", "utf-8-sig", "cp949", "euc-kr"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="ignore")


def _extract_pdf_text(raw: bytes) -> str:
    reader = PdfReader(BytesIO(raw))
    pages: list[str] = []
    for page in reader.pages:
        pages.append(page.extract_text() or "")
    return "\n".join(pages)


def _extract_ppt_text(raw: bytes) -> str:
    presentation = Presentation(BytesIO(raw))
    lines: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if text:
                lines.append(text)
    return "\n".join(lines)


def _normalize_whitespace(text: str) -> str:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _chunk_text(text: str, chunk_size: int = 900, overlap: int = 140) -> list[str]:
    cleaned = _normalize_whitespace(text)
    if not cleaned:
        return []
    chunks: list[str] = []
    start = 0
    length = len(cleaned)
    while start < length and len(chunks) < MAX_CHUNKS_PER_SESSION:
        end = min(start + chunk_size, length)
        if end < length:
            boundary = cleaned.rfind("\n", start, end)
            if boundary > start + 250:
                end = boundary
        piece = cleaned[start:end].strip()
        if piece:
            chunks.append(piece)
        if end >= length:
            break
        start = max(0, end - overlap)
    return chunks


def _tokenize(text: str) -> set[str]:
    return {tok for tok in re.findall(r"[0-9A-Za-z가-힣]{2,}", text.lower())}


def _extract_rag_save_payload(message: str) -> str | None:
    # Examples:
    # - "A는 B야. RAG 정보로 추가해줘"
    # - "A는 B야 rag에 저장해줘"
    # - "A는 B야\n* rag 저장"
    original = (message or "").strip()
    if not original:
        return None

    # New explicit command format:
    #   <memo text>
    #   * rag 저장
    bullet_command = re.search(r"(?im)^\s*[*-]\s*rag\s*저장\s*$", original)
    if bullet_command:
        payload = original[:bullet_command.start()].strip(" \n\t:：-")
        if payload:
            return payload

    # Also support inline suffix command: "<memo text> * rag 저장"
    inline_suffix = re.search(r"(?i)\*\s*rag\s*저장\s*$", original)
    if inline_suffix:
        payload = original[:inline_suffix.start()].strip(" \n\t:：-")
        if payload:
            return payload

    lowered = message.lower()
    trigger_patterns = [
        r"\brag\s*정보\s*로?\s*추가\s*해\s*줘\b",
        r"\brag\s*에\s*추가\s*해\s*줘\b",
        r"\brag\s*에\s*저장\s*해\s*줘\b",
        r"\brag\s*저장\s*해\s*줘\b",
        r"\brag\s*정보\s*로?\s*영구\s*저장\s*해\s*줘\b",
        r"\brag\s*영구\s*저장\s*해\s*줘\b",
    ]
    for pattern in trigger_patterns:
        match = re.search(pattern, lowered)
        if not match:
            continue
        # Use original casing/string for saved payload
        original_start = match.start()
        payload = message[:original_start].strip(" \n\t:：-")
        if payload:
            return payload
    return None


def _is_rag_save_only_command(message: str) -> bool:
    text = (message or "").strip()
    if not text:
        return False
    return bool(
        re.fullmatch(r"(?is)[*-]?\s*rag\s*저장\s*(해\s*줘)?\s*", text)
    )


def _extract_manual_note_title(note_text: str) -> tuple[str | None, str]:
    """
    Parse optional title markup from manual RAG note text.
    Supported form:
    - </제목>
    """
    text = (note_text or "").strip()
    if not text:
        return None, ""

    match = re.search(r"</\s*([^<>\n]{1,200})\s*>", text)
    if not match:
        return None, text

    raw_title = re.sub(r"\s+", " ", match.group(1)).strip(" \n\t:：-")
    title = raw_title[:200] if raw_title else None
    body = (text[:match.start()] + text[match.end():]).strip(" \n\t:：-")
    return title, body


def _latest_user_note_from_history(history: List[ChatMessage]) -> str | None:
    for msg in reversed(history or []):
        if msg.role != "user":
            continue
        content = (msg.content or "").strip()
        if not content:
            continue
        if _is_rag_save_only_command(content):
            continue
        if _extract_rag_save_payload(content) is not None:
            # Skip historical "save command" turns.
            continue
        return content
    return None


def _append_rag_manual_note(session_id: str, note_text: str, title: str | None = None) -> int:
    chunks = _chunk_text(note_text, chunk_size=700, overlap=80)
    if not chunks:
        return 0
    docs = RAG_STORE.setdefault(session_id, [])
    file_name = title if title else "[수동저장] 사용자 메모"
    docs.append(
        {
            "doc_id": str(uuid.uuid4()),
            "file_name": file_name,
            "chunks": chunks,
            "tokens": [_tokenize(c) for c in chunks],
        }
    )

    total_chunks = sum(len(d["chunks"]) for d in docs)
    while total_chunks > MAX_CHUNKS_PER_SESSION and docs:
        removed = docs.pop(0)
        total_chunks -= len(removed["chunks"])
    _save_rag_store()
    return len(chunks)


def _build_rag_context(session_id: str, query: str, top_k: int = 4) -> str:
    docs = RAG_STORE.get(session_id) or []
    if not docs:
        return ""
    q_tokens = _tokenize(query)
    if not q_tokens:
        return ""

    scored: list[tuple[float, str, str, int]] = []
    for doc in docs:
        file_name = doc["file_name"]
        for idx, chunk in enumerate(doc["chunks"]):
            tokens = doc["tokens"][idx]
            overlap = len(q_tokens.intersection(tokens))
            if overlap == 0:
                continue
            score = overlap / max(1, len(q_tokens))
            scored.append((score, file_name, chunk, idx + 1))

    if not scored:
        return ""

    scored.sort(key=lambda x: x[0], reverse=True)
    selected = scored[:top_k]
    lines = [
        "아래는 사용자가 업로드한 문서에서 검색된 참고 문맥입니다. 근거가 있을 때만 활용하세요:",
    ]
    for _, file_name, chunk, index in selected:
        lines.append(f"[출처: {file_name} / 청크 {index}]")
        lines.append(chunk)
    return "\n\n".join(lines)


def _build_rag_context_all_sessions(query: str, top_k: int = 4) -> str:
    q_tokens = _tokenize(query)
    if not q_tokens:
        return ""

    scored: list[tuple[float, str, str, int, str]] = []
    for session_id, docs in RAG_STORE.items():
        for doc in docs:
            file_name = doc["file_name"]
            for idx, chunk in enumerate(doc["chunks"]):
                tokens = doc["tokens"][idx]
                overlap = len(q_tokens.intersection(tokens))
                if overlap == 0:
                    continue
                score = overlap / max(1, len(q_tokens))
                scored.append((score, file_name, chunk, idx + 1, session_id))

    if not scored:
        return ""

    scored.sort(key=lambda x: x[0], reverse=True)
    selected = scored[:top_k]
    lines = [
        "아래는 저장된 전체 RAG 문서에서 검색된 참고 문맥입니다. 근거가 있을 때만 활용하세요:",
    ]
    for _, file_name, chunk, index, session_id in selected:
        lines.append(f"[출처: {file_name} / 청크 {index} / 세션 {session_id[:8]}]")
        lines.append(chunk)
    return "\n\n".join(lines)


def _find_doc_index(session_id: str, doc_id: str) -> int:
    docs = RAG_STORE.get(session_id) or []
    for idx, doc in enumerate(docs):
        if str(doc.get("doc_id", "")).strip() == doc_id:
            return idx
    return -1


def _find_doc_index_by_name(session_id: str, file_name: str) -> int:
    docs = RAG_STORE.get(session_id) or []
    for idx, doc in enumerate(docs):
        if str(doc.get("file_name", "")).strip() == file_name:
            return idx
    return -1


@app.on_event("startup")
async def _startup_load_rag_store():
    _load_rag_store()


async def _resolve_model_name(requested_model: str) -> str:
    data = await ollama_request("GET", "/api/tags", timeout=30.0)
    models = [m.get("name") for m in data.get("models", []) if m.get("name")]
    if not models:
        return requested_model

    if requested_model in models:
        return requested_model

    latest_candidate = f"{requested_model}:latest"
    if latest_candidate in models:
        return latest_candidate

    stripped = requested_model.split(":", 1)[0]
    for model in models:
        if model.split(":", 1)[0] == stripped:
            return model

    return requested_model


def _is_vision_model(model_name: str) -> bool:
    lower = model_name.lower()
    return any(key in lower for key in VISION_MODEL_HINTS)


async def _resolve_vision_model(preferred_model: str) -> str | None:
    data = await ollama_request("GET", "/api/tags", timeout=30.0)
    models = [m.get("name") for m in data.get("models", []) if m.get("name")]
    if not models:
        return None

    # Keep requested model if it is already a vision model and installed.
    if preferred_model in models and _is_vision_model(preferred_model):
        return preferred_model

    preferred_base = preferred_model.split(":", 1)[0]
    for model in models:
        model_base = model.split(":", 1)[0]
        if model_base == preferred_base and _is_vision_model(model):
            return model

    for model in models:
        if _is_vision_model(model):
            return model

    return None


@app.post("/api/chat", response_model=ChatResponse)
async def chat(request: ChatRequest):
    user_message = request.message.strip()
    if not user_message:
        raise HTTPException(status_code=400, detail="메시지는 비어 있을 수 없습니다.")

    # Deterministic datetime response to prevent model from answering stale dates.
    if _is_datetime_query(user_message):
        return ChatResponse(response=_format_datetime_answer(), model="system-clock")

    # Command mode: "RAG 정보로 추가해줘" -> persist note in RAG store.
    manual_note = _extract_rag_save_payload(user_message)
    if manual_note is None and _is_rag_save_only_command(user_message):
        manual_note = _latest_user_note_from_history(request.history)
    if manual_note is not None:
        session_id = (request.rag_session_id or "").strip() or "global-default"
        note_title, note_body = _extract_manual_note_title(manual_note)
        added_chunks = _append_rag_manual_note(session_id, note_body, note_title)
        if added_chunks == 0:
            raise HTTPException(status_code=400, detail="저장할 RAG 텍스트를 찾지 못했습니다.")
        saved_title = note_title if note_title else "[수동저장] 사용자 메모"
        return ChatResponse(
            response=(
                "요청한 내용을 RAG에 영구 저장했습니다.\n"
                f"- 제목: {saved_title}\n"
                f"- 저장 청크 수: {added_chunks}\n"
                "이후 질문부터 이 정보를 참고해 답변합니다."
            ),
            model=request.model,
        )

    resolved_model = await _resolve_model_name(request.model)
    if request.image_base64:
        vision_model = await _resolve_vision_model(resolved_model)
        if not vision_model:
            raise HTTPException(
                status_code=400,
                detail=(
                    "이미지 해석용 비전 모델이 설치되어 있지 않습니다. "
                    "다음 명령으로 설치 후 다시 시도해주세요: ollama pull gemma3:latest "
                    "(또는 ollama pull llava:latest)"
                ),
            )
        resolved_model = vision_model

    messages = [{"role": msg.role, "content": msg.content} for msg in request.history]
    if not messages or messages[0].get("role") != "system":
        messages.insert(0, {"role": "system", "content": _current_date_system_message()})

    if request.use_rag and request.rag_session_id:
        rag_context = _build_rag_context(request.rag_session_id, user_message)
        if not rag_context:
            rag_context = _build_rag_context_all_sessions(user_message)
        if rag_context:
            messages.insert(
                1,
                {
                    "role": "system",
                    "content": (
                        "문서 기반 답변 규칙:\n"
                        "- 아래 문맥을 우선 참고해 답하세요.\n"
                        "- 문맥에 없는 내용은 추측하지 말고 모른다고 말하세요.\n\n"
                        + rag_context
                    ),
                },
            )

    user_turn: dict = {"role": "user", "content": user_message}
    if request.image_base64:
        user_turn["images"] = [request.image_base64]
    messages.append(user_turn)

    try:
        data = await ollama_request(
            "POST",
            "/api/chat",
            payload={
                "model": resolved_model,
                "messages": messages,
                "stream": False,
                "options": {"temperature": 0.7, "top_p": 0.9},
            },
        )
        content = data.get("message", {}).get("content", "")
    except HTTPException as exc:
        # Ollama 구버전(또는 일부 환경)에서 /api/chat 미지원 시 /api/generate로 폴백
        if exc.status_code != 404:
            raise
        fallback_payload = {
            "model": resolved_model,
            "prompt": "System: " + _current_date_system_message() + "\n" + _history_to_prompt(request.history, user_message),
            "stream": False,
            "options": {"temperature": 0.7, "top_p": 0.9},
        }
        if request.image_base64:
            fallback_payload["images"] = [request.image_base64]

        fallback = await ollama_request(
            "POST",
            "/api/generate",
            payload=fallback_payload,
        )
        content = fallback.get("response", "")

    if not content:
        content = "응답이 없습니다."

    return ChatResponse(response=content, model=resolved_model)


@app.post("/api/rag/upload", response_model=RagUploadResponse)
async def rag_upload(session_id: str = Form(...), file: UploadFile = File(...)):
    file_name = (file.filename or "").strip()
    if not file_name:
        raise HTTPException(status_code=400, detail="파일 이름이 없습니다.")

    raw = await file.read()
    if not raw:
        raise HTTPException(status_code=400, detail="빈 파일은 업로드할 수 없습니다.")
    if len(raw) > MAX_UPLOAD_BYTES:
        raise HTTPException(status_code=400, detail="파일 크기는 50MB 이하만 지원합니다.")

    ext = os.path.splitext(file_name)[1].lower()
    try:
        if ext == ".txt":
            extracted = _decode_text_file(raw)
        elif ext == ".pdf":
            extracted = _extract_pdf_text(raw)
        elif ext in (".pptx", ".ppt"):
            extracted = _extract_ppt_text(raw)
        else:
            raise HTTPException(status_code=400, detail="지원하지 않는 파일 형식입니다. (pdf, ppt, pptx, txt)")
    except HTTPException:
        raise
    except Exception:
        if ext == ".ppt":
            raise HTTPException(
                status_code=400,
                detail="PPT(구형 바이너리) 파싱에 실패했습니다. PPTX로 저장 후 업로드해주세요.",
            )
        raise HTTPException(status_code=400, detail="파일 텍스트 추출에 실패했습니다.")

    chunks = _chunk_text(extracted)
    if not chunks:
        raise HTTPException(status_code=400, detail="문서에서 추출된 텍스트가 없습니다.")

    normalized_session = session_id.strip() or str(uuid.uuid4())
    doc = {
        "doc_id": str(uuid.uuid4()),
        "file_name": file_name,
        "chunks": chunks,
        "tokens": [_tokenize(c) for c in chunks],
    }
    docs = RAG_STORE.setdefault(normalized_session, [])
    docs.append(doc)

    # Limit memory for long-running process.
    total_chunks = sum(len(d["chunks"]) for d in docs)
    while total_chunks > MAX_CHUNKS_PER_SESSION and docs:
        removed = docs.pop(0)
        total_chunks -= len(removed["chunks"])

    _save_rag_store()

    return RagUploadResponse(
        session_id=normalized_session,
        doc_id=doc["doc_id"],
        file_name=file_name,
        chunks=len(chunks),
    )


@app.post("/api/rag/clear")
async def rag_clear(request: RagSessionClearRequest):
    RAG_STORE.pop(request.session_id, None)
    _save_rag_store()
    return {"status": "ok", "session_id": request.session_id}


@app.get("/api/rag/list")
async def rag_list(session_id: str):
    docs = RAG_STORE.get(session_id) or []
    items = []
    for idx, doc in enumerate(docs):
        items.append(
            {
                "doc_index": idx,
                "doc_id": doc.get("doc_id", ""),
                "file_name": doc["file_name"],
                "chunks": len(doc["chunks"]),
            }
        )
    return {"session_id": session_id, "files": items}


@app.get("/api/rag/list_all")
async def rag_list_all():
    items: list[dict] = []
    for session_id, docs in RAG_STORE.items():
        for idx, doc in enumerate(docs):
            items.append(
                {
                    "session_id": session_id,
                    "doc_index": idx,
                    "doc_id": doc.get("doc_id", ""),
                    "file_name": doc.get("file_name", ""),
                    "chunks": len(doc.get("chunks", [])),
                }
            )
    return {"files": items}


@app.post("/api/rag/rename")
async def rag_rename(request: RagRenameRequest):
    session_id = (request.session_id or "").strip()
    doc_id = (request.doc_id or "").strip()
    doc_index = request.doc_index
    old_file_name = (request.old_file_name or "").strip()
    new_name = (request.file_name or "").strip()
    if not session_id:
        raise HTTPException(status_code=400, detail="session_id는 필수입니다.")
    if not doc_id and doc_index is None and not old_file_name:
        raise HTTPException(status_code=400, detail="doc_id/doc_index/old_file_name 중 하나가 필요합니다.")
    if not new_name:
        raise HTTPException(status_code=400, detail="파일 이름은 비어 있을 수 없습니다.")

    idx = _find_doc_index(session_id, doc_id) if doc_id else -1
    if idx < 0 and doc_index is not None:
        docs = RAG_STORE.get(session_id) or []
        if 0 <= int(doc_index) < len(docs):
            idx = int(doc_index)
    if idx < 0 and old_file_name:
        idx = _find_doc_index_by_name(session_id, old_file_name)
    if idx < 0:
        raise HTTPException(status_code=404, detail="문서를 찾을 수 없습니다.")

    target_doc = RAG_STORE[session_id][idx]
    if not str(target_doc.get("doc_id", "")).strip():
        target_doc["doc_id"] = str(uuid.uuid4())
    target_doc["file_name"] = new_name
    _save_rag_store()
    return {
        "status": "ok",
        "session_id": session_id,
        "doc_id": target_doc.get("doc_id", ""),
        "file_name": new_name,
    }


@app.post("/api/rag/delete")
async def rag_delete(request: RagDeleteRequest):
    session_id = (request.session_id or "").strip()
    doc_id = (request.doc_id or "").strip()
    doc_index = request.doc_index
    old_file_name = (request.old_file_name or "").strip()
    if not session_id:
        raise HTTPException(status_code=400, detail="session_id는 필수입니다.")
    if not doc_id and doc_index is None and not old_file_name:
        raise HTTPException(status_code=400, detail="doc_id/doc_index/old_file_name 중 하나가 필요합니다.")

    idx = _find_doc_index(session_id, doc_id) if doc_id else -1
    if idx < 0 and doc_index is not None:
        docs = RAG_STORE.get(session_id) or []
        if 0 <= int(doc_index) < len(docs):
            idx = int(doc_index)
    if idx < 0 and old_file_name:
        idx = _find_doc_index_by_name(session_id, old_file_name)
    if idx < 0:
        raise HTTPException(status_code=404, detail="문서를 찾을 수 없습니다.")

    docs = RAG_STORE.get(session_id) or []
    removed = docs.pop(idx)
    if not docs:
        RAG_STORE.pop(session_id, None)
    _save_rag_store()
    return {
        "status": "ok",
        "session_id": session_id,
        "doc_id": removed.get("doc_id", ""),
        "file_name": removed.get("file_name", ""),
    }


@app.get("/api/rag/document")
async def rag_document_get(
    session_id: str,
    doc_id: str | None = None,
    doc_index: int | None = None,
    old_file_name: str | None = None,
):
    normalized_session = (session_id or "").strip()
    normalized_doc_id = (doc_id or "").strip()
    normalized_old_name = (old_file_name or "").strip()
    if not normalized_session:
        raise HTTPException(status_code=400, detail="session_id는 필수입니다.")
    if not normalized_doc_id and doc_index is None and not normalized_old_name:
        raise HTTPException(status_code=400, detail="doc_id/doc_index/old_file_name 중 하나가 필요합니다.")

    idx = _find_doc_index(normalized_session, normalized_doc_id) if normalized_doc_id else -1
    if idx < 0 and doc_index is not None:
        docs = RAG_STORE.get(normalized_session) or []
        if 0 <= int(doc_index) < len(docs):
            idx = int(doc_index)
    if idx < 0 and normalized_old_name:
        idx = _find_doc_index_by_name(normalized_session, normalized_old_name)
    if idx < 0:
        raise HTTPException(status_code=404, detail="문서를 찾을 수 없습니다.")

    doc = RAG_STORE[normalized_session][idx]
    content = "\n\n".join([c for c in doc.get("chunks", []) if isinstance(c, str)]).strip()
    return {
        "status": "ok",
        "session_id": normalized_session,
        "doc_index": idx,
        "doc_id": doc.get("doc_id", ""),
        "file_name": doc.get("file_name", ""),
        "content": content,
        "chunks": len(doc.get("chunks", [])),
    }


@app.post("/api/rag/update_content")
async def rag_document_update(request: RagDocumentUpdateRequest):
    session_id = (request.session_id or "").strip()
    doc_id = (request.doc_id or "").strip()
    doc_index = request.doc_index
    old_file_name = (request.old_file_name or "").strip()
    next_content = (request.content or "").strip()
    if not session_id:
        raise HTTPException(status_code=400, detail="session_id는 필수입니다.")
    if not doc_id and doc_index is None and not old_file_name:
        raise HTTPException(status_code=400, detail="doc_id/doc_index/old_file_name 중 하나가 필요합니다.")
    if not next_content:
        raise HTTPException(status_code=400, detail="내용은 비어 있을 수 없습니다.")

    idx = _find_doc_index(session_id, doc_id) if doc_id else -1
    if idx < 0 and doc_index is not None:
        docs = RAG_STORE.get(session_id) or []
        if 0 <= int(doc_index) < len(docs):
            idx = int(doc_index)
    if idx < 0 and old_file_name:
        idx = _find_doc_index_by_name(session_id, old_file_name)
    if idx < 0:
        raise HTTPException(status_code=404, detail="문서를 찾을 수 없습니다.")

    chunks = _chunk_text(next_content)
    if not chunks:
        raise HTTPException(status_code=400, detail="유효한 텍스트를 찾을 수 없습니다.")

    target_doc = RAG_STORE[session_id][idx]
    if not str(target_doc.get("doc_id", "")).strip():
        target_doc["doc_id"] = str(uuid.uuid4())
    target_doc["chunks"] = chunks
    target_doc["tokens"] = [_tokenize(c) for c in chunks]
    _save_rag_store()
    return {
        "status": "ok",
        "session_id": session_id,
        "doc_id": target_doc.get("doc_id", ""),
        "file_name": target_doc.get("file_name", ""),
        "chunks": len(chunks),
    }


@app.get("/api/models")
async def list_models():
    data = await ollama_request("GET", "/api/tags", timeout=30.0)
    models = data.get("models", [])

    model_names = sorted([m.get("name") for m in models if m.get("name")])
    return {"models": model_names, "default": DEFAULT_MODEL}


@app.get("/api/health")
async def health_check():
    try:
        await ollama_request("GET", "/api/version", timeout=10.0)
        return {"status": "ok", "ollama": "connected", "ollama_base_url": OLLAMA_BASE_URL}
    except HTTPException as exc:
        return {
            "status": "degraded",
            "ollama": "disconnected",
            "detail": exc.detail,
            "ollama_base_url": OLLAMA_BASE_URL,
        }


@app.get("/")
async def root():
    return {"message": "NOS LLM Chat API", "status": "running"}


if __name__ == "__main__":
    import uvicorn

    def _is_port_available(host: str, port: int) -> bool:
        with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as sock:
            sock.setsockopt(socket.SOL_SOCKET, socket.SO_REUSEADDR, 1)
            try:
                sock.bind((host, port))
                return True
            except OSError:
                return False

    def _resolve_backend_port(preferred_port: int, host: str = "0.0.0.0") -> int:
        # Prefer deterministic fallback that frontend already probes.
        candidates: list[int] = [preferred_port]
        if preferred_port != 8010:
            candidates.append(8010)
        candidates.extend(range(preferred_port + 1, preferred_port + 21))
        for candidate in candidates:
            if _is_port_available(host, candidate):
                return candidate
        return preferred_port

    requested_port = int(os.getenv("BACKEND_PORT", "8000"))
    resolved_port = _resolve_backend_port(requested_port)
    if resolved_port != requested_port:
        print(
            f"[NOS] Port {requested_port} is already in use. "
            f"Starting backend on port {resolved_port} instead."
        )
    uvicorn.run(app, host="0.0.0.0", port=resolved_port)

